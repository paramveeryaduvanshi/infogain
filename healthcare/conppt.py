from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re

# Read the markdown file
with open(r"c:\Users\Paramveer Singh\OneDrive\Project\infogain\healthcare\PRESENTATION.md", "r", encoding="utf-8") as f:
    content = f.read()

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
TITLE_COLOR = RGBColor(0, 51, 102)  # Dark blue
ACCENT_COLOR = RGBColor(0, 153, 204)  # Light blue
TEXT_COLOR = RGBColor(50, 50, 50)  # Dark gray

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 245, 250)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = TITLE_COLOR
    title_para.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(32)
        subtitle_para.font.color.rgb = ACCENT_COLOR
        subtitle_para.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_lines):
    """Add a content slide with bullet points"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = TITLE_COLOR
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.6), Inches(5.8))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, line in enumerate(content_lines):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = line
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(6)
        p.space_after = Pt(6)

# Slide 1: Title slide
add_title_slide(prs, "Healthcare AI Chatbot", "Clinical Intelligence Platform")

# Slide 2: Executive Summary
add_content_slide(prs, "Executive Summary", [
    "🔒 Privacy-First: All processing done locally",
    "🤖 AI-Powered: LLM-driven clinical insights",
    "📊 Data-Driven: Multi-source health data integration",
    "⚡ Real-Time: Instant clinical recommendations",
    "",
    "A Django-based healthcare platform using local LLMs for intelligent patient data analysis"
])

# Slide 3: Architecture
add_content_slide(prs, "System Architecture", [
    "User Interface Layer: Web-based chat frontend",
    "Backend Engine: Django with REST API",
    "AI Core: LLM-driven SQL generation (Ollama)",
    "Data Layer: SQLite database with multiple datasets",
    "Quality Assurance: RAGAS-based evaluation",
    "Deployment: Containerizable, scalable infrastructure"
])

# Slide 4: Technology Stack
add_content_slide(prs, "Technology Stack", [
    "Frontend: HTML5, CSS, JavaScript",
    "Backend: Django 5.2.12, Python 3.11",
    "Database: SQLite3 (local)",
    "LLM: Ollama with llama3.2 model",
    "AI Libraries: LangChain, LangGraph, RAGAS",
    "Data Processing: Pandas, NumPy"
])

# Slide 5: Core Components
add_content_slide(prs, "Core Components", [
    "Data Ingestion: Excel → SQLite with null handling",
    "Query Agent: Natural language → SQL conversion",
    "Health Analysis: Patient data → Clinical insights",
    "Frontend: Interactive web chat interface",
    "Evaluation: RAGAS metrics for quality assurance",
    "Prompt Engine: Role-based structured output"
])

# Slide 6: Data Pipeline
add_content_slide(prs, "Data Pipeline", [
    "1. Load Excel files (Health Dataset 1 & 2)",
    "2. Clean & validate data (null values → 0)",
    "3. Store in SQLite database",
    "4. LLM generates SQL from natural language",
    "5. Query execution & data retrieval",
    "6. Format output as {PatientX: {attributes}}"
])

# Slide 7: Query Processing Flow
add_content_slide(prs, "Intelligent Query Processing", [
    "User Query: 'What is patient 1's health status?'",
    "Schema Discovery: Identify available tables/columns",
    "LLM Translation: Convert to SQL using Ollama",
    "Query Execution: Run SQL on local database",
    "Data Transformation: Nested dict format",
    "Clinical Analysis: Generate recommendations"
])

# Slide 8: Key Features
add_content_slide(prs, "Key Features", [
    "✓ Natural Language Queries → SQL generation",
    "✓ Patient Data Retrieval in structured format",
    "✓ Clinical Analysis & Recommendations",
    "✓ Local LLM Inference (no cloud dependency)",
    "✓ Quality Metrics via RAGAS",
    "✓ Null Value Handling (automatic fillna)"
])

# Slide 9: Use Cases
add_content_slide(prs, "Use Cases", [
    "Patient Health Queries: Real-time status checks",
    "Symptom Analysis: Diagnosis support",
    "Clinical Decision Support: Personalized plans",
    "Quality Assurance: Response accuracy metrics",
    "Epidemiological Studies: Population health analysis",
    "Training & Research: Clinical decision patterns"
])

# Slide 10: Privacy & Security
add_content_slide(prs, "Privacy & Security", [
    "✓ 100% Local Processing - No cloud APIs",
    "✓ Complete Data Control",
    "✓ HIPAA-Compliant Infrastructure Ready",
    "✓ No Third-Party Data Sharing",
    "✓ SQLite Database Stored Locally",
    "✓ Offline Capability"
])

# Slide 11: Performance Metrics
add_content_slide(prs, "Performance Metrics", [
    "Query Response Time: <5 seconds",
    "Data Processing: <1 second",
    "Throughput: Multi-query support",
    "Latency: Sub-second (local)",
    "Accuracy: RAGAS-measured",
    "Scalability: Supports concurrent requests"
])

# Slide 12: Competitive Advantages
add_content_slide(prs, "Competitive Advantages", [
    "Local-First: No vendor lock-in",
    "Cost-Effective: One-time setup vs per-API pricing",
    "Low Latency: Network-independent processing",
    "Full Customization: Complete source access",
    "Privacy: Patient data never leaves facility",
    "Scalable: From single-user to enterprise"
])

# Slide 13: Implementation Status
add_content_slide(prs, "Implementation Status", [
    "✓ Core LLM-SQL agent (Complete)",
    "✓ Patient data querying (Complete)",
    "✓ Clinical recommendation engine (Complete)",
    "✓ RAGAS evaluation framework (Complete)",
    "✓ Web interface (Complete)",
    "✓ Data ingestion pipeline (Complete)"
])

# Slide 14: Roadmap
add_content_slide(prs, "Development Roadmap", [
    "Phase 1 (Current): Core features ✓",
    "Phase 2 (Planned): Multi-turn conversations, advanced filtering",
    "Phase 3 (Planned): Federated learning, mobile app, EHR integration",
    "",
    "Enhanced Features:",
    "• Historical tracking • Medication interaction checking • Predictive analytics"
])

# Slide 15: Project Structure
add_content_slide(prs, "Project Structure", [
    "healthcare/",
    "├── chatagent.py - SQL generation & execution",
    "├── views.py - API endpoints",
    "├── prompt.py - LLM instruction templates",
    "├── evaluationRagas.py - Quality metrics",
    "└── chatbot.html - Web interface",
    ""
])

# Slide 16: Setup & Deployment
add_content_slide(prs, "Quick Start", [
    "1. Activate virtual environment",
    "2. Install dependencies (Django, LangChain, Ollama)",
    "3. Run: python healthcare/ingestdata.py",
    "4. Run: python manage.py runserver",
    "5. Access: http://localhost:8000/chatbot/",
    ""
])

# Slide 17: Technical Innovation
add_content_slide(prs, "Technical Innovation", [
    "LLM-Driven SQL: Dynamic query generation",
    "Nested Data Format: {PatientX: {attributes}}",
    "Structured Prompting: Hallucination prevention",
    "Local-First AI: Privacy + Performance",
    "Schema Discovery: Auto-adapting queries",
    "Quality Metrics: RAGAS integration"
])

# Slide 18: Next Steps
add_content_slide(prs, "Next Steps", [
    "✓ User testing with clinical staff",
    "✓ Expand patient datasets",
    "✓ Multi-turn conversation support",
    "✓ Advanced analytics dashboard",
    "✓ EHR system integration",
    "✓ Production deployment"
])

# Slide 19: Closing
add_title_slide(prs, "Thank You!", "Healthcare AI Chatbot Platform\nPrivacy-First Clinical Intelligence")

# Save presentation
output_path = r"c:\Users\Paramveer Singh\OneDrive\Project\infogain\healthcare\PRESENTATION.pptx"
prs.save(output_path)
print(f"✓ PowerPoint presentation created: {output_path}")